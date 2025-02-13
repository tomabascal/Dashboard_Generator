# Importing the required libraries
import streamlit as st
import pandas as pd
import pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
import time
import zipfile
import io
import shutil
from datetime import datetime
import re
import subprocess
from openpyxl import load_workbook


# Function to convert PPTX to PDF
def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Convert a PPTX file to PDF using LibreOffice (works on Streamlit Cloud)."""
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf",
                       pptx_path, "--outdir", os.path.dirname(pdf_path)], check=True)
    except Exception as e:
        print(f"Error converting {pptx_path} to PDF: {e}")


# Function to create a ZIP file with all the generated PPTX files
def create_zip_of_presentations(folder_path):
    """Create a ZIP file with all the generated PPTX files in the folder."""
    zip_buffer = io.BytesIO()

    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zipf:
        for file in os.listdir(folder_path):
            file_path = os.path.join(folder_path, file)
            if file.endswith(".pptx"):  # Avoid including template and Excel
                zipf.write(file_path, arcname=file)

    zip_buffer.seek(0)
    return zip_buffer


# Function to get the file name based on the selected columns
def get_filename_from_selection(row, selected_columns):
    """Generate the file name based on the selected columns."""
    file_name_parts = [str(int(row[col])) if isinstance(row[col], float) and row[col].is_integer() else str(row[col])
                       for col in selected_columns if col in row]
    return "_".join(file_name_parts)


# Function to update the text of text boxes in the PPTX file
def update_text_of_textbox(presentation, column_letter, new_text, wb, sheet_name, cell_coordinate):
    """Find and replace text within text boxes that have the format {A}, {B}, etc."""
    pattern = rf"\{{{column_letter}\}}"  # Regular expression to find "{A}", "{B}", etc.

    # Format the text according to the Excel format
    formatted_text = format_cell_value(
        new_text, wb, sheet_name, cell_coordinate)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                if re.search(pattern, shape.text):  # Search for pattern in text
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = re.sub(pattern, str(
                                formatted_text), run.text)  # Replace


# Function to process the files and generate the reports
def process_files(ppt_file, excel_file, search_option, start_row, end_row, store_ids, selected_columns, output_format):
    """Generate reports in PPTX or PDF format on Streamlit Cloud with estimated time notification."""

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    folder_name = f"Presentations_{timestamp}"
    os.makedirs(folder_name, exist_ok=True)

    temp_folder = "temp_files"
    os.makedirs(temp_folder, exist_ok=True)

    ppt_template_path = os.path.join(temp_folder, ppt_file.name)
    excel_file_path = os.path.join(temp_folder, excel_file.name)

    with open(ppt_template_path, "wb") as f:
        f.write(ppt_file.getbuffer())
    with open(excel_file_path, "wb") as f:
        f.write(excel_file.getbuffer())

    try:
        wb = load_workbook(excel_file_path, data_only=True)
        df1 = pd.read_excel(excel_file_path, sheet_name=0,
                            dtype={'Store ID': str})
        sheet_name = wb.sheetnames[0]  # Get the name of the first sheet
    except PermissionError as e:
        st.error(f"Error reading Excel file: {e}")
        return

    if search_option == 'rows':
        df_selected = df1.iloc[start_row-2:end_row-1]
    elif search_option == 'store_id':
        store_id_list = [store_id.strip() for store_id in store_ids.split(',')]
        df_selected = df1[df1.iloc[:, 0].astype(str).isin(store_id_list)]
    else:
        df_selected = pd.DataFrame()

    total_files = len(df_selected)
    if total_files == 0:
        st.error("⚠️ No files to generate. Check the filters.")
        return

    # 🔹 Estimated time notification based on the chosen format
    estimated_time = total_files * (5 if output_format == "PDF" else 1)
    st.info(f"⏳ Estimated time: ~{estimated_time} seconds")

    progress_bar = st.progress(0)
    progress_text = st.empty()

    current_file = 0
    start_time = time.time()

    for index, row in df_selected.iterrows():
        process_row(ppt_template_path, row, df1, index,
                    selected_columns, folder_name, output_format, wb, sheet_name)
        current_file += 1
        progress = current_file / total_files
        progress_bar.progress(progress)
        elapsed_time = time.time() - start_time
        progress_text.write(
            f"📄 Generating {current_file}/{total_files} ({output_format}) - Elapsed time: {int(elapsed_time)}s")

    # Create a ZIP with the files in the selected format
    zip_path = f"{folder_name}.zip"
    shutil.make_archive(zip_path.replace(".zip", ""), 'zip', folder_name)

    with open(zip_path, "rb") as zip_file:
        st.download_button(
            label=f"📥 Download {total_files} reports ({output_format})",
            data=zip_file,
            file_name=f"{folder_name}.zip",
            mime="application/zip"
        )

    progress_text.write(
        f"✅ All reports have been generated in {output_format} format! Total time: {int(time.time() - start_time)}s")


# Function to process a row and generate a PPTX or PDF file
def process_row(presentation_path, row, df1, index, selected_columns, output_folder, output_format, wb, sheet_name):
    """Process a row and generate a PPTX or PDF file on Streamlit Cloud."""
    presentation = pptx.Presentation(presentation_path)

    for col_idx, col_name in enumerate(row.index):
        column_letter = chr(65 + col_idx)
        # Adjust for Excel's 1-based index
        cell_coordinate = f"{column_letter}{index + 2}"
        update_text_of_textbox(presentation, column_letter,
                               row[col_name], wb, sheet_name, cell_coordinate)

    file_name = get_filename_from_selection(row, selected_columns)
    pptx_path = os.path.join(output_folder, f"{file_name}.pptx")

    # Save as PPTX
    presentation.save(pptx_path)

    # If the user chooses PDF, convert the file
    if output_format == "PDF":
        pdf_path = os.path.join(output_folder, f"{file_name}.pdf")
        convert_pptx_to_pdf(pptx_path, pdf_path)
        # Delete the original PPTX to only keep the PDF
        os.remove(pptx_path)


# Function to format Excel cell values based on their type
def format_cell_value(value, wb, sheet_name, cell_coordinate=None):
    """Formats and rounds the cell value based on its type and format in Excel."""
    if value is None:
        return ""

    if isinstance(value, (int, float)) and cell_coordinate:
        ws = wb[sheet_name]
        cell = ws[cell_coordinate]
        cell_format = cell.number_format

        # Clean strange characters from the format (e.g., \#,##0\ "€")
        # Removed ',' from the regex
        cleaned_format = re.sub(r'[^\d.%€$£]', '', cell_format)

        # Identify the currency symbol if it exists
        currency_symbol = next(
            (symbol for symbol in ["€", "$", "£"] if symbol in cleaned_format), "")

        if currency_symbol:
            # Round to 1 decimal and remove the .0 if it is an integer
            rounded_value = round(value, 1)
            return f"{rounded_value:.1f}".rstrip('0').rstrip('.') + f" {currency_symbol}"
        elif "%" in cleaned_format:
            # Round percentage to 1 decimal, but never show .0
            percentage = round(value * 100, 1)
            if percentage.is_integer():  # If the percentage is an integer
                return f"{int(percentage)}%"  # Do not show decimals
            else:
                return f"{percentage:.1f}%"  # Show decimals
        else:
            # Round normal number to 1 decimal and remove the .0 if it is an integer
            rounded_value = round(value, 1)
            return f"{rounded_value:.1f}".rstrip('0').rstrip('.')

    elif isinstance(value, datetime):
        return value.strftime("%d-%m-%Y")  # Date format

    return str(value)


# ========= 💡 Styles to improve design =========
st.markdown("""
    <style>
    div.stButton > button {
        width: 100%;
        height: 50px;
        font-size: 16px;
        border-radius: 10px;
        font-weight: bold;
    }
    </style>
""", unsafe_allow_html=True)

# ========= Title =========
st.title("Shopfully Target By Store Reporting Tool")

# Option to choose output format
st.markdown("### **Select Output Format** \n**`PPTX Format Recommended`**")
output_format = st.radio("Choose the file format:", ["PPTX", "PDF"])

# Warning message if the user chooses PDF
if output_format == "PDF":
    st.warning(
        "⚠️ Converting to PDF may take extra time. Large batches of presentations might take several minutes.")

# ========= 📂 Upload files with improved format =========
st.markdown(
    "**Upload PPTX Template**  \n*(Text Box format that will be edited -> {Column Letter} For Example: `{A}`)*", unsafe_allow_html=True)
ppt_template = st.file_uploader("", type=["pptx"])

st.write("")  # Spacing

st.markdown(
    "**Upload Excel File**  \n*(Column A must be `Store ID`)*", unsafe_allow_html=True)
data_file = st.file_uploader("", type=["xlsx"])


# ========= 🔍 Improved buttons for "Search by" =========
st.markdown("### **Search by:**")  # Bold and larger title
col1, col2 = st.columns(2)  # Two columns to align buttons in a grid

# Initialize the state variable for filter selection
if "search_option" not in st.session_state:
    st.session_state.search_option = "rows"  # Default value

# Button 1 - Search by Rows
with col1:
    if st.button("🔢 Rows", use_container_width=True):
        st.session_state.search_option = "rows"

# Button 2 - Search by Store ID
with col2:
    if st.button("🔍 Store ID", use_container_width=True):
        st.session_state.search_option = "store_id"

# Show the selected option
st.markdown(f"**Selected: `{st.session_state.search_option}`**")


# ========= 🔢 Inputs to define the search range =========
start_row, end_row, store_ids = None, None, None

if st.session_state.search_option == "rows":
    start_row = st.number_input("Start Row", min_value=1, step=1)
    end_row = st.number_input("End Row", min_value=1, step=1)

elif st.session_state.search_option == "store_id":
    store_ids = st.text_input("Enter Store IDs (comma-separated)")


# ========= 📝 Column selection for file name =========
if data_file is not None:
    # Read the first sheet of the Excel file
    df = pd.read_excel(data_file, sheet_name=0, dtype={'Store ID': str})
    column_names = df.columns.tolist()

    selected_columns = st.multiselect(
        "📂 Select and order the columns for the file name:",
        column_names,
        default=column_names[:1]
    )

    st.write("🔹 Example file name:", get_filename_from_selection(
        df.iloc[0], selected_columns))


# ========= 🚀 Processing button =========
if st.button("Process"):
    if ppt_template and data_file:
        process_files(ppt_template, data_file, st.session_state.search_option,
                      start_row, end_row, store_ids, selected_columns, output_format)
    else:
        st.error("Please upload both files before processing.")
