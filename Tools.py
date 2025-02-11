# Load the necessary libraries
import streamlit as st
import pandas as pd
import pptx
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
import time
import zipfile
import io
import shutil
from datetime import datetime
import re
import subprocess
import openpyxl
from openpyxl import load_workbook


# Function to convert a PPTX file to PDF using LibreOffice (works on Streamlit Cloud)
def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Converts a PPTX file to PDF on Linux using LibreOffice (works on Streamlit Cloud)."""
    try:
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf",
                       pptx_path, "--outdir", os.path.dirname(pdf_path)], check=True)
    except Exception as e:
        print(f"Error converting {pptx_path} to PDF: {e}")


# Function to create a ZIP file with all the generated PPTX files
def create_zip_of_presentations(folder_path):
    """Creates a ZIP file with all the generated PPTX files in the folder."""
    zip_buffer = io.BytesIO()

    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zipf:
        for file in os.listdir(folder_path):
            file_path = os.path.join(folder_path, file)
            if file.endswith(".pptx"):  # Evitamos incluir plantilla y Excel
                zipf.write(file_path, arcname=file)

    zip_buffer.seek(0)
    return zip_buffer


# Function to generate the file name based on the selected columns
def get_filename_from_selection(row, selected_columns):
    """Generates the file name based on the selected columns."""
    file_name_parts = [str(row[col]) for col in selected_columns if col in row]
    return "_".join(file_name_parts)


# Function to update the text of a textbox in a PPTX presentation
def update_text_of_textbox(presentation, column_letter, new_text):
    """Searches for and replaces text inside text boxes with the format {A}, {B}, etc., while preserving the PPTX formatting."""
    pattern = rf"\{{{column_letter}\}}"

    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                if re.search(pattern, shape.text):
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = re.sub(pattern, str(new_text), run.text)


def process_files(ppt_file, excel_file, search_option, start_row, end_row, store_ids, selected_columns, output_format):
    """Generates reports in PPTX or PDF format on Streamlit Cloud while preserving the Excel formatting."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    folder_name = f"Presentations_{timestamp}"
    os.makedirs(folder_name, exist_ok=True)

    temp_folder = "temp_files"
    os.makedirs(temp_folder, exist_ok=True)

    # Save temporary files.
    ppt_template_path = os.path.join(temp_folder, ppt_file.name)
    excel_file_path = os.path.join(temp_folder, excel_file.name)

    with open(ppt_template_path, "wb") as f:
        f.write(ppt_file.getbuffer())
    with open(excel_file_path, "wb") as f:
        f.write(excel_file.getbuffer())

    # Leer el archivo Excel con pandas para filtrar datos
    try:
        with pd.ExcelFile(excel_file_path) as xls:
            df1 = pd.read_excel(xls, sheet_name=0)
    except PermissionError as e:
        st.error(f"Error reading Excel file: {e}")
        return

    if search_option == 'rows':
        df_selected = df1.iloc[start_row:end_row + 1]
    elif search_option == 'store_id':
        store_id_list = [store_id.strip() for store_id in store_ids.split(',')]
        df_selected = df1[df1.iloc[:, 0].astype(str).isin(store_id_list)]
    else:
        df_selected = pd.DataFrame()

    total_files = len(df_selected)
    if total_files == 0:
        st.error("⚠️ No data found. Verify filters.")
        return

    st.info(f"⏳ Estimated time: ~{total_files} seconds")

    
    progress_bar = st.progress(0)
    progress_text = st.empty()

    start_time = time.time()
    for i, (_, row) in enumerate(df_selected.iterrows()):  # Enumerate to get a sequential counter
        process_row(ppt_template_path, row, excel_file_path, i, selected_columns, folder_name, output_format)
        progress_bar.progress((i + 1) / total_files)  # Use the sequential counter
        progress_text.write(f"📄 Processing {i + 1}/{total_files}")

    # Create ZIP with the generated files
    zip_path = f"{folder_name}.zip"
    shutil.make_archive(zip_path.replace(".zip", ""), 'zip', folder_name)

    with open(zip_path, "rb") as zip_file:
        st.download_button(
            label=f"📥 Download {total_files} reports ({output_format})",
            data=zip_file,
            file_name=f"{folder_name}.zip",
            mime="application/zip"
        )

    progress_text.write(f"✅ Finished! Total time: {int(time.time() - start_time)}s")


# Function to process a row and generate a PPTX or PDF file while preserving the original Excel formatting
def process_row(presentation_path, row, excel_file_path, index, selected_columns, output_folder, output_format):
    """Processes a row and generates a PPTX or PDF file while preserving the original Excel formatting."""    
    # Load the PowerPoint presentation
    presentation = pptx.Presentation(presentation_path)

    # Load the Excel file with openpyxl to read the formats
    wb = load_workbook(excel_file_path, data_only=True)
    ws = wb.active  # Read first sheet of the workbook

    for col_idx, col_name in enumerate(row.index):
        column_letter = chr(65 + col_idx)  # Convert column index to letter (A, B, C, ...)
        excel_cell = ws[f"{column_letter}{index + 2}"]  # Index + 2 because Excel is 1-indexed

        # Get the formatted value
        formatted_text = format_cell_value(excel_cell, wb, ws.title)
        
        # Replace the text in the presentation
        update_text_of_textbox(presentation, column_letter, formatted_text)

    # Generate the file name based on the selected columns
    file_name = get_filename_from_selection(row, selected_columns)
    pptx_path = os.path.join(output_folder, f"{file_name}.pptx")

    # Save the PPTX file
    presentation.save(pptx_path)

    # Convert to PDF if necessary
    if output_format == "PDF":
        pdf_path = os.path.join(output_folder, f"{file_name}.pdf")
        convert_pptx_to_pdf(pptx_path, pdf_path)
        os.remove(pptx_path)  # Delete the PPTX after conversion


# Function to format Excel cell values based on their type
def format_cell_value(cell, wb, sheet_name):
    """Formats and rounds the cell value based on its type and format in Excel."""
    if cell is None or cell.value is None:
        return ""
    
    value = cell.value
    if isinstance(value, (int, float)):
        ws = wb[sheet_name]
        cell_format = ws[cell.coordinate].number_format

        # Clean strange characters from the format (e.g., \#,##0\ "€")
        cleaned_format = re.sub(r'[^\d.,%€$£]', '', cell_format)  

        # Identify the currency symbol if it exists
        currency_symbol = next((symbol for symbol in ["€", "$", "£"] if symbol in cleaned_format), "")

        if currency_symbol:
            # Round to 1 decimal and remove the .0 if it is an integer
            rounded_value = round(value, 1)
            return f"{rounded_value:,.1f}".rstrip('0').rstrip('.') + f" {currency_symbol}"
        elif "%" in cleaned_format:
            # Round percentage to 1 decimal, but never show .0
            percentage = round(value * 100, 1)
            if percentage.is_integer():  # If the percentage is an integer
                return f"{int(percentage)}%"  # Do not show decimals
            else:
                return f"{percentage:.1f}%"  # Show decimals
        else:
            # Round normal number to 1 decimal and remove the .0 if it is an integer
            rounded_value = round(value, 1)
            return f"{rounded_value:,.1f}".rstrip('0').rstrip('.')  # Round to 1 decimal

    elif isinstance(value, datetime):
        return value.strftime("%d-%m-%Y")  # Date format

    return str(value)

# ========= 💡 Styles to enhance the design =========
st.markdown("""
    <style>
    div.stButton > button {
        width: 100%;
        height: 50px;
        font-size: 16px;
        border-radius: 10px;
        font-weight: bold;
    }
    </style>
""", unsafe_allow_html=True)

# ========= Title =========
st.title("Shopfully Target By Store Reporting Tool")

# Option to choose the output format
st.markdown("### **Select Output Format**")
output_format = st.radio("Choose the file format:", ["PPTX `(recommended)`", "PDF"])

# Warning message if the user chooses PDF
if output_format == "PDF":
    st.warning(
        "⚠️ Converting to PDF may take extra time. Large batches of presentations might take several minutes.")

# ========= 📂 Enhanced file upload =========
st.markdown(
    "**Upload PPTX Template**  \n*(Text Box format that will be edited -> {Column Letter} For Example: **`{A}`**)*", unsafe_allow_html=True)
ppt_template = st.file_uploader("", type=["pptx"])

st.write("")  # Space between file uploaders

st.markdown(
    "**Upload Excel File**  \n*(Column A must be **Store ID**)*", unsafe_allow_html=True)
data_file = st.file_uploader("", type=["xlsx"])

# ========= 🔍 # Buttons for "Search by" =========
st.markdown("### **Search by:**")  
col1, col2 = st.columns(2)  

if "search_option" not in st.session_state:
    st.session_state.search_option = "rows"  # Default search option

# Button 1 - Search by Rows
with col1:
    if st.button("🔢 Rows", use_container_width=True):
        st.session_state.search_option = "rows"

# Button 2 - Search by Store ID
with col2:
    if st.button("🔍 Store ID", use_container_width=True):
        st.session_state.search_option = "store_id"

# Display the selected option
st.markdown(f"**Selected: {st.session_state.search_option}**")

# ========= 🔢 Inputs to define search range  =========
start_row, end_row, store_ids = None, None, None

if st.session_state.search_option == "rows":
    start_row = st.number_input("Start Row", min_value=1, step=1)
    end_row = st.number_input("End Row", min_value=1, step=1)

elif st.session_state.search_option == "store_id":
    store_ids = st.text_input("Enter Store IDs (comma-separated)")

# ========= 📝 Column selector for file name =========
if data_file is not None:
    df = pd.read_excel(data_file, sheet_name=0)
    column_names = df.columns.tolist()

    # Apply filters based on the selected option
    if st.session_state.search_option == "rows" and start_row is not None and end_row is not None:
        df = df.iloc[start_row:end_row]  # Adjust for zero-based index

    elif st.session_state.search_option == "store_id" and store_ids:
        store_ids_list = [int(id.strip()) for id in store_ids.split(',')]
        df = df[df['store_id'].isin(store_ids_list)]

    selected_columns = st.multiselect(
        "📂 Select and order the columns for the file name:",
        column_names,
        default=column_names[:1]
    )

    def get_filename_from_selection(row, selected_columns):
        """Generates the file name based on the selected columns."""
        file_name_parts = [str(row[col])
                           for col in selected_columns if col in row]
        return "_".join(file_name_parts)

if len(df) > 1:
    st.write("🔹 Example file name:", get_filename_from_selection(
        df.iloc[1], selected_columns))
else:
    st.warning("The DataFrame does not have enough rows to display an example file name.")


# ========= 🚀 Process Button =========
if st.button("Process"):
    if ppt_template and data_file:
        process_files(ppt_template, data_file, st.session_state.search_option,
                      start_row, end_row, store_ids, selected_columns, output_format)
    else:
        st.error("Please upload both files before processing.")